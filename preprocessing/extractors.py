"""
Document text extractors for various file formats.
Each extractor returns a list of (page/section, text) tuples along with metadata.
"""

import os
from datetime import datetime, timezone
from typing import Any

import pdfplumber
import docx
import openpyxl
from pptx import Presentation
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup


def extract_metadata_base(filepath: str) -> dict[str, Any]:
    """Extract common metadata from a file path."""
    filename = os.path.basename(filepath)
    ext = os.path.splitext(filename)[1].lower().lstrip(".")
    return {
        "filename": filename,
        "filetype": ext,
        "title": os.path.splitext(filename)[0],
        "author": None,
        "page_count": None,
        "processed_at": datetime.now(timezone.utc).isoformat(),
    }


def extract_pdf(filepath: str) -> tuple[dict[str, Any], list[tuple[int, str]]]:
    """
    Extract text and metadata from a PDF file using pdfplumber.

    Uses tuned x_tolerance for accurate word boundary detection,
    which avoids the mid-word space artifacts common with PyPDF2
    on justified text.

    Returns:
        Tuple of (metadata dict, list of (page_number, text) tuples).
    """
    metadata = extract_metadata_base(filepath)
    pages: list[tuple[int, str]] = []

    with pdfplumber.open(filepath) as pdf:
        # Extract metadata from PDF info dict
        info = pdf.metadata or {}
        metadata["title"] = info.get("Title") or metadata["title"]
        metadata["author"] = info.get("Author")
        metadata["page_count"] = len(pdf.pages)

        for i, page in enumerate(pdf.pages):
            text = page.extract_text(
                x_tolerance=3,  # Horizontal tolerance for grouping chars into words
                y_tolerance=3,  # Vertical tolerance for grouping chars into lines
            ) or ""
            if text.strip():
                pages.append((i + 1, text))

    return metadata, pages


def extract_docx(filepath: str) -> tuple[dict[str, Any], list[tuple[int, str]]]:
    """
    Extract text and metadata from a Word document.

    Returns:
        Tuple of (metadata dict, list of (section_number, text) tuples).
    """
    metadata = extract_metadata_base(filepath)

    doc = docx.Document(filepath)
    core = doc.core_properties
    metadata["title"] = core.title or metadata["title"]
    metadata["author"] = core.author

    paragraphs: list[tuple[int, str]] = []
    section_num = 1
    current_section: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Split on headings to create logical sections
        if para.style and para.style.name.startswith("Heading"):
            if current_section:
                paragraphs.append((section_num, "\n".join(current_section)))
                section_num += 1
                current_section = []
        current_section.append(text)

    if current_section:
        paragraphs.append((section_num, "\n".join(current_section)))

    metadata["page_count"] = section_num

    return metadata, paragraphs


def _format_cell(cell) -> str:
    """Format a single cell value for readable text output."""
    if cell is None:
        return ""
    if isinstance(cell, datetime):
        # Strip time component if it's midnight (date-only values)
        if cell.hour == 0 and cell.minute == 0 and cell.second == 0:
            return cell.strftime("%Y-%m-%d")
        return cell.isoformat()
    if isinstance(cell, bool):
        return "Yes" if cell else "No"
    if isinstance(cell, float):
        # Clean up float display: 3.0 → "3", 3.85 → "3.85"
        if cell == int(cell):
            return str(int(cell))
        return f"{cell:.4g}"
    return str(cell).strip()


def extract_xlsx(filepath: str) -> tuple[dict[str, Any], list[tuple[int, str]]]:
    """
    Extract text and metadata from an Excel file.

    Produces self-contained row-group chunks where each chunk carries
    column headers so it remains meaningful in isolation.

    Format per chunk:
        [Sheet: SheetName]
        Row 1: Column1=Value1 | Column2=Value2 | ...
        Row 2: Column1=Value1 | Column2=Value2 | ...

    Returns:
        Tuple of (metadata dict, list of (sheet_number, text) tuples).
    """
    metadata = extract_metadata_base(filepath)

    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    sections: list[tuple[int, str]] = []

    ROWS_PER_CHUNK = 15  # Group rows into chunks of this size

    for sheet_idx, sheet_name in enumerate(wb.sheetnames):
        ws = wb[sheet_name]

        # Collect all rows
        all_rows = list(ws.iter_rows(values_only=True))
        if not all_rows:
            continue

        # Detect header row (first non-empty row)
        headers: list[str] = []
        data_start = 0
        for row_idx, row in enumerate(all_rows):
            non_empty = [c for c in row if c is not None]
            if non_empty:
                headers = [str(c).strip() if c is not None else f"Col{i+1}"
                           for i, c in enumerate(row)]
                data_start = row_idx + 1
                break

        if not headers:
            continue

        # Collect data rows
        data_rows: list[list[str]] = []
        for row in all_rows[data_start:]:
            formatted = [_format_cell(c) for c in row]
            # Skip fully empty rows
            if any(v for v in formatted):
                data_rows.append(formatted)

        if not data_rows:
            # Sheet only has headers, output them anyway
            header_line = " | ".join(headers)
            sections.append((
                sheet_idx + 1,
                f"[Sheet: {sheet_name}]\nColumns: {header_line}\n(No data rows)"
            ))
            continue

        # Split data rows into groups for chunking
        for group_start in range(0, len(data_rows), ROWS_PER_CHUNK):
            group = data_rows[group_start:group_start + ROWS_PER_CHUNK]
            lines: list[str] = [f"[Sheet: {sheet_name}]"]

            for row_offset, row_values in enumerate(group):
                row_num = data_start + group_start + row_offset + 1
                # Format as key=value pairs
                pairs: list[str] = []
                for header, value in zip(headers, row_values):
                    if value:  # Skip empty cells
                        pairs.append(f"{header}={value}")
                if pairs:
                    lines.append(f"Row {row_num}: {' | '.join(pairs)}")

            if len(lines) > 1:  # More than just the sheet header
                sections.append((sheet_idx + 1, "\n".join(lines)))

    metadata["page_count"] = len(wb.sheetnames)
    wb.close()

    return metadata, sections


def extract_pptx(filepath: str) -> tuple[dict[str, Any], list[tuple[int, str]]]:
    """
    Extract text and metadata from a PowerPoint file.

    Returns:
        Tuple of (metadata dict, list of (slide_number, text) tuples).
    """
    metadata = extract_metadata_base(filepath)

    prs = Presentation(filepath)
    core = prs.core_properties
    metadata["title"] = core.title or metadata["title"]
    metadata["author"] = core.author
    metadata["page_count"] = len(prs.slides)

    slides: list[tuple[int, str]] = []

    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append((i + 1, "\n".join(texts)))

    return metadata, slides


def extract_epub(filepath: str) -> tuple[dict[str, Any], list[tuple[int, str]]]:
    """
    Extract text and metadata from an EPUB file.

    Returns:
        Tuple of (metadata dict, list of (chapter_number, text) tuples).
    """
    metadata = extract_metadata_base(filepath)

    book = epub.read_epub(filepath)

    # Extract metadata
    title = book.get_metadata("DC", "title")
    if title:
        metadata["title"] = title[0][0]

    creator = book.get_metadata("DC", "creator")
    if creator:
        metadata["author"] = creator[0][0]

    chapters: list[tuple[int, str]] = []
    chapter_num = 1

    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        content = item.get_content()
        soup = BeautifulSoup(content, "html.parser")
        text = soup.get_text(separator="\n", strip=True)

        if text.strip():
            chapters.append((chapter_num, text))
            chapter_num += 1

    metadata["page_count"] = len(chapters)

    return metadata, chapters


# Registry mapping file extensions to their extractor functions
EXTRACTORS: dict[str, callable] = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "xlsx": extract_xlsx,
    "pptx": extract_pptx,
    "epub": extract_epub,
}


def get_extractor(filepath: str):
    """Get the appropriate extractor function for a given file."""
    ext = os.path.splitext(filepath)[1].lower().lstrip(".")
    extractor = EXTRACTORS.get(ext)
    if extractor is None:
        raise ValueError(
            f"Unsupported file format: .{ext}. "
            f"Supported formats: {', '.join(EXTRACTORS.keys())}"
        )
    return extractor
